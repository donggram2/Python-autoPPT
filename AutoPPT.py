import collections
import collections.abc
import os
import copy
import pandas

from datetime import datetime
from pptx import Presentation
from pptx.util import Cm

os.chdir("/Users/wondongha/PycharmProjects/autoPPT/")
cwd = os.getcwd()
print(cwd)
images = [file for file in os.listdir(os.path.join(cwd, "images")) if os.path.isfile(os.path.join(cwd, "images", file))]

# 대체 이미지 파일 경로
default_image = os.path.join(cwd, "default_image.png")


def convert_date_format(date):
    if pandas.notnull(date):  # null 값이 아닌 경우에만 변환
        weekday = ['월', '화', '수', '목', '금', '토', '일']
        weekday_idx = date.weekday()  # 요일 인덱스 (0: 월요일, 1: 화요일, ...)
        weekday_korean = weekday[weekday_idx]  # 요일 한글 표현
        return date.strftime("'%y.%m.%d") + f"({weekday_korean})"
    return date
def select_table_by_text(slide, text):
    for x in slide.shapes:
        if x.has_table and x.table.cell(0, 0).text == text:
            return x.table
    print('요청한 Shape를 찾을 수 없습니다.')


def copy_slide(prs, index):
    template = prs.slides[index]
    try:
        blank_slide_layout = prs.slide_layouts.get_by_name(BLANK_LAYOUT_NAME)
    except:
        blank_slide_layout = prs.slide_layouts[0]
    copied_slide = prs.slides.add_slide(blank_slide_layout)

    for shape in template.shapes:
        elem = shape.element
        new_elem = copy.deepcopy(elem)
        copied_slide.shapes._spTree.insert_element_before(new_elem, 'p:extLst')
    return copied_slide


def update_text(text_frame, new_text):
    p = text_frame.paragraphs[0]
    run = p.runs[0]
    run.text = new_text


# 경로 변수 설정
xlsx_path = os.path.join(cwd, "test.xlsx")
pptx_template = os.path.join(cwd, "Auto_PPT_Format1.pptx")
pptx_result = os.path.join(cwd, "Auto_PPT_Result.pptx")

# 상수 변수 정의
SLIDE_INDEX = 0
BLANK_LAYOUT_NAME = '빈화면'

df = pandas.read_excel(xlsx_path, sheet_name="Sheet1", skiprows=2, converters={"통행일자": convert_date_format})
prs = Presentation(pptx_template)
slide = prs.slides[SLIDE_INDEX]

for i, r in df.iterrows():
    copied_slide = copy_slide(prs, SLIDE_INDEX)
    table = select_table_by_text(copied_slide, '위반 내용')

    update_text(table.cell(1, 1).text_frame, f"{r['위반속도(km)']:,} km/h")
    update_text(table.cell(2, 1).text_frame, str(r['장소']))
    update_text(table.cell(0, 3).text_frame, str(r['위반 업체 및 부서명']) + "\n" + str(r['차량번호']))
    update_text(table.cell(1, 3).text_frame, str(r['통행일자']) + "\n" + f"{r['통행시간']}")

    if len(images) > 0 and any(str(r['차량번호']) in x for x in images):
        image = [x for x in images if str(r['차량번호']) in x][0]
        image_path = os.path.join(cwd, "images", image)
    else:
        image_path = default_image

    inserted_logo = copied_slide.shapes.add_picture(image_path, left=Cm(1.38), top=Cm(3.81))
    inserted_logo.width = Cm(13.78)
    inserted_logo.height = Cm(8.38)

prs.save(pptx_result)
